"""
Document Text Extraction Utilities
Extracts text content from various document formats for RAG pipeline
Includes OCR support for image-based PDFs
"""
import io
import os
from typing import Tuple
import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import fitz  # PyMuPDF
from PIL import Image, ImageEnhance, ImageFilter
import pytesseract

# Configure Tesseract path from environment variable (Windows support)
if os.getenv("TESSERACT_CMD"):
    pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_CMD")


def preprocess_image_for_ocr(image: Image.Image) -> Image.Image:
    """
    Preprocess image to improve OCR accuracy.
    
    Args:
        image: PIL Image object
    
    Returns:
        Preprocessed PIL Image
    """
    # Convert to grayscale
    image = image.convert('L')
    
    # Increase contrast
    enhancer = ImageEnhance.Contrast(image)
    image = enhancer.enhance(2.0)
    
    # Increase sharpness
    enhancer = ImageEnhance.Sharpness(image)
    image = enhancer.enhance(2.0)
    
    # Apply slight denoising (median filter)
    image = image.filter(ImageFilter.MedianFilter(size=3))
    
    # Optional: Increase brightness slightly
    enhancer = ImageEnhance.Brightness(image)
    image = enhancer.enhance(1.1)
    
    return image


def extract_text_from_pdf_with_ocr(file_content: bytes) -> Tuple[str, dict]:
    """
    Extract text from image-based PDF using OCR (PyMuPDF + Tesseract).
    Uses enhanced preprocessing for better accuracy on scanned documents.
    
    Args:
        file_content: Binary content of PDF file
    
    Returns:
        Tuple of (extracted_text, metadata)
    """
    try:
        print("   🔍 Using OCR for image-based PDF...")
        
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(stream=file_content, filetype="pdf")
        text_content = []
        num_pages = len(pdf_document)
        ocr_pages = 0
        total_chars = 0
        
        # Tesseract configurations to try (in order of preference)
        # PSM modes: 3=auto, 6=single uniform block, 11=sparse text, 1=auto with OSD
        tesseract_configs = [
            '--oem 3 --psm 3',  # Default: Auto page segmentation with OEM (LSTM)
            '--oem 3 --psm 6',  # Single uniform block of text
            '--oem 3 --psm 4',  # Single column of text
        ]
        
        for page_num in range(num_pages):
            page = pdf_document[page_num]
            
            # Convert page to high-resolution image (3x zoom for better OCR)
            # Higher DPI = better OCR for scanned documents
            pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))
            img_data = pix.tobytes("png")
            
            # Open as PIL Image
            image = Image.open(io.BytesIO(img_data))
            
            # Preprocess image for better OCR
            processed_image = preprocess_image_for_ocr(image)
            
            # Perform OCR with multiple config attempts
            page_text = ""
            best_text = ""
            
            for config in tesseract_configs:
                try:
                    page_text = pytesseract.image_to_string(
                        processed_image, 
                        lang='eng',
                        config=config
                    )
                    
                    # Keep the result with most content
                    if len(page_text.strip()) > len(best_text.strip()):
                        best_text = page_text
                    
                    # If we got good content, stop trying other configs
                    if len(page_text.strip()) > 100:
                        break
                        
                except Exception as ocr_error:
                    print(f"   ⚠️ OCR config '{config}' failed for page {page_num + 1}: {str(ocr_error)}")
                    continue
            
            # Use the best result
            if best_text.strip():
                # Clean up common OCR artifacts
                cleaned_text = best_text.strip()
                
                # Remove common watermarks/artifacts
                lines = [line for line in cleaned_text.split('\n') 
                        if line.strip() and 
                        not line.strip().lower() in ['scanned with camscanner', 'camscanner']]
                
                cleaned_text = '\n'.join(lines)
                
                if cleaned_text.strip():
                    text_content.append(f"[Page {page_num + 1} - OCR]\n{cleaned_text}")
                    ocr_pages += 1
                    total_chars += len(cleaned_text)
                    print(f"   📄 Page {page_num + 1}: Extracted {len(cleaned_text)} characters")
            else:
                print(f"   ⚠️ No text extracted from page {page_num + 1}")
        
        pdf_document.close()
        
        full_text = "\n\n".join(text_content)
        
        metadata = {
            "num_pages": num_pages,
            "ocr_pages": ocr_pages,
            "total_chars": total_chars,
            "extraction_method": "ocr_enhanced",
            "has_content": len(full_text.strip()) > 0
        }
        
        print(f"   ✅ OCR completed: {ocr_pages}/{num_pages} pages, {total_chars} total characters extracted")
        
        return full_text, metadata
        
    except Exception as e:
        raise Exception(f"Failed to extract text from PDF with OCR: {str(e)}")


def extract_text_from_pdf(file_content: bytes) -> Tuple[str, dict]:
    """
    Extract text from PDF file. Falls back to OCR if regular extraction fails.
    
    Args:
        file_content: Binary content of PDF file
    
    Returns:
        Tuple of (extracted_text, metadata)
    """
    try:
        # First, try regular text extraction
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text_content = []
        num_pages = len(pdf_reader.pages)
        
        for page_num, page in enumerate(pdf_reader.pages, 1):
            page_text = page.extract_text()
            if page_text.strip():
                text_content.append(f"[Page {page_num}]\n{page_text}")
        
        full_text = "\n\n".join(text_content)
        
        # If insufficient text found (likely image-based PDF), use OCR
        if len(full_text.strip()) < 50:
            print("   ⚠️ Insufficient text found with regular extraction, trying OCR...")
            return extract_text_from_pdf_with_ocr(file_content)
        
        metadata = {
            "num_pages": num_pages,
            "extraction_method": "text",
            "has_content": len(full_text.strip()) > 0
        }
        
        return full_text, metadata
        
    except Exception as e:
        # If regular extraction fails completely, try OCR as last resort
        try:
            print(f"   ⚠️ Regular extraction failed: {str(e)}, trying OCR...")
            return extract_text_from_pdf_with_ocr(file_content)
        except Exception as ocr_error:
            raise Exception(f"Failed to extract text from PDF (both regular and OCR): {str(ocr_error)}")


def extract_text_from_docx(file_content: bytes) -> Tuple[str, dict]:
    """
    Extract text from Word document (.docx).
    
    Args:
        file_content: Binary content of Word file
    
    Returns:
        Tuple of (extracted_text, metadata)
    """
    try:
        docx_file = io.BytesIO(file_content)
        doc = Document(docx_file)
        
        text_content = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_content.append(row_text)
        
        full_text = "\n".join(text_content)
        
        metadata = {
            "num_paragraphs": len(doc.paragraphs),
            "num_tables": len(doc.tables),
            "has_content": len(full_text.strip()) > 0
        }
        
        return full_text, metadata
        
    except Exception as e:
        raise Exception(f"Failed to extract text from DOCX: {str(e)}")


def extract_text_from_xlsx(file_content: bytes) -> Tuple[str, dict]:
    """
    Extract text from Excel spreadsheet (.xlsx).
    
    Args:
        file_content: Binary content of Excel file
    
    Returns:
        Tuple of (extracted_text, metadata)
    """
    try:
        xlsx_file = io.BytesIO(file_content)
        workbook = load_workbook(xlsx_file, data_only=True)
        
        text_content = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_content.append(f"[Sheet: {sheet_name}]")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text_content.append(row_text)
        
        full_text = "\n".join(text_content)
        
        metadata = {
            "num_sheets": len(workbook.sheetnames),
            "sheet_names": workbook.sheetnames,
            "has_content": len(full_text.strip()) > 0
        }
        
        return full_text, metadata
        
    except Exception as e:
        raise Exception(f"Failed to extract text from XLSX: {str(e)}")


def extract_text_from_pptx(file_content: bytes) -> Tuple[str, dict]:
    """
    Extract text from PowerPoint presentation (.pptx).
    
    Args:
        file_content: Binary content of PowerPoint file
    
    Returns:
        Tuple of (extracted_text, metadata)
    """
    try:
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        
        text_content = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text_content.append(f"[Slide {slide_num}]")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
        
        full_text = "\n".join(text_content)
        
        metadata = {
            "num_slides": len(presentation.slides),
            "has_content": len(full_text.strip()) > 0
        }
        
        return full_text, metadata
        
    except Exception as e:
        raise Exception(f"Failed to extract text from PPTX: {str(e)}")


def extract_text_from_txt(file_content: bytes) -> Tuple[str, dict]:
    """
    Extract text from plain text file (.txt).
    
    Args:
        file_content: Binary content of text file
    
    Returns:
        Tuple of (extracted_text, metadata)
    """
    try:
        # Try UTF-8 first, fallback to latin-1
        try:
            text = file_content.decode('utf-8')
        except UnicodeDecodeError:
            text = file_content.decode('latin-1', errors='ignore')
        
        metadata = {
            "encoding": "utf-8" if text == file_content.decode('utf-8', errors='ignore') else "latin-1",
            "num_lines": len(text.splitlines()),
            "has_content": len(text.strip()) > 0
        }
        
        return text, metadata
        
    except Exception as e:
        raise Exception(f"Failed to extract text from TXT: {str(e)}")


def extract_text_from_document(file_content: bytes, mime_type: str, file_name: str) -> Tuple[str, dict]:
    """
    Extract text from document based on MIME type.
    
    Args:
        file_content: Binary content of the file
        mime_type: MIME type of the file
        file_name: Name of the file (for extension fallback)
    
    Returns:
        Tuple of (extracted_text, metadata)
    """
    # Map MIME types to extractors
    mime_extractors = {
        'application/pdf': extract_text_from_pdf,
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': extract_text_from_docx,
        'application/msword': extract_text_from_docx,
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': extract_text_from_xlsx,
        'application/vnd.ms-excel': extract_text_from_xlsx,
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': extract_text_from_pptx,
        'application/vnd.ms-powerpoint': extract_text_from_pptx,
        'text/plain': extract_text_from_txt,
    }
    
    # Try by MIME type first
    if mime_type in mime_extractors:
        return mime_extractors[mime_type](file_content)
    
    # Fallback: Try by file extension
    file_ext = file_name.lower().split('.')[-1] if '.' in file_name else ''
    ext_extractors = {
        'pdf': extract_text_from_pdf,
        'docx': extract_text_from_docx,
        'doc': extract_text_from_docx,
        'xlsx': extract_text_from_xlsx,
        'xls': extract_text_from_xlsx,
        'pptx': extract_text_from_pptx,
        'ppt': extract_text_from_pptx,
        'txt': extract_text_from_txt,
    }
    
    if file_ext in ext_extractors:
        return ext_extractors[file_ext](file_content)
    
    raise Exception(f"Unsupported document type: {mime_type} / .{file_ext}")
